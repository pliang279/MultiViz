from pptx import Presentation, enum
from pptx.util import Inches, Pt
import json

# CLIP Flickr30k Analysis - with Double Grad

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

from structured_framework.misc.flickr30k_clip_target_ids import *

id_to_tids = {
    50: instance_text_target_ids_50,
    100: instance_text_target_ids_100,
    150: instance_text_target_ids_150,
    200: instance_text_target_ids_200,
    500: instance_text_target_ids_500,
    250: instance_text_target_ids_250,
    300: instance_text_target_ids_300,
    400: instance_text_target_ids_400,
    600: instance_text_target_ids_600,
    700: instance_text_target_ids_700,
    800: instance_text_target_ids_800,
    900: instance_text_target_ids_900,
    1000: instance_text_target_ids_1000,
    350: instance_text_target_ids_350,
    450: instance_text_target_ids_450,
    550: instance_text_target_ids_550,
    650: instance_text_target_ids_650,
    750: instance_text_target_ids_750,
    850: instance_text_target_ids_850,
    950: instance_text_target_ids_950
}



for instance_idx, tid_dict in id_to_tids.items():
    # Read Key to Logits JSON
    with open(f"clip_key_to_logits/key_to_logits-box-acc-{instance_idx}.json") as f:
        key_to_logits = json.loads(f.read())
    for key, value in tid_dict.items(): 
        orig_img_path = f"structured_framework/visuals/flickr30k-clip-{instance_idx}-{0}-image.png"
        orig_text = open(f"structured_framework/visuals/flickr30k-clip-{instance_idx}-{0}-text.txt").read()
        doublegrad_img_path = f"structured_framework/visuals/flickr30k-clip-{key}-doublegrad.png"
        new_img_path = f"structured_framework/visuals/flickr30k-clip-{key}-new_box_img_with_boxes.jpg"
        new_text = open(f"structured_framework/visuals/flickr30k-clip-{key}-new_text.txt").read()
        random_img_path =  f"structured_framework/visuals/flickr30k-clip-{key}-random_box_img.jpg"
        gt_img_path = f"structured_framework/visuals/flickr30k-clip-{key}-gt_img.jpg"
        slide = prs.slides.add_slide(blank_slide_layout)


        # "original_logits": 10.043411,
        # "doublegrad_logits": 3.0287437,
        # "random_drop_logits": 8.995966,
        # "ground_truth_logits": 8.610129

        orig_lgts = key_to_logits[str(instance_idx)][key]["original_logits"]
        dg_lgts = key_to_logits[str(instance_idx)][key]["doublegrad_box_logits"]
        rd_lgts = key_to_logits[str(instance_idx)][key]["random_box_logits"]
        gt_lgts = key_to_logits[str(instance_idx)][key]["ground_truth_logits"]


        # Original Image
        left = top = Inches(0.5)
        height = Inches(2.5)
        pic = slide.shapes.add_picture(orig_img_path, left, top, height=height)
        
        # Original Text
        left = Inches(1)
        top = Inches(3)
        height = Inches(2)
        width = Inches(3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        p = tf.add_paragraph()
        if len(orig_text.split()) > 15:
            third = len(orig_text.split()) // 3
            text = ' '.join(orig_text.split()[:third] + ["\n"]+ orig_text.split()[third:third*2] + ["\n"] + orig_text.split()[third*2:])
        else:
            mid = len(orig_text.split()) // 2
            text = ' '.join(orig_text.split()[:mid] + ["\n"] + orig_text.split()[mid:])
        p.text = text
        p.word_wrap = True
        p.font.size = Pt(10)

        left = Inches(1)
        top = Inches(3.5)
        height = Inches(0.5)
        width = Inches(3)
        txBox2 = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox2.text_frame

        double_grad_text = value["text"]

        if len(double_grad_text.split()) > 15:
            third = len(double_grad_text.split()) // 3
            text = ' '.join(double_grad_text.split()[:third] + ["\n"]+ double_grad_text.split()[third:third*2] + ["\n"] + double_grad_text.split()[third*2:])
        elif len(double_grad_text.split()) > 10:
            mid = len(orig_text.split()) // 2
            text = ' '.join(orig_text.split()[:mid] + ["\n"] + orig_text.split()[mid:])
        else:
            text = double_grad_text

        p = tf.add_paragraph()
        p.text = f"Double Grad Text: {text} \n Orig Logits: {orig_lgts:.3f}"
        p.word_wrap = True
        p.font.size = Pt(10)

        # DoubleGrad Image
        left = Inches(5)
        top = Inches(0.5)
        height = Inches(2.5)
        pic = slide.shapes.add_picture(doublegrad_img_path, left, top, height=height)

        # New Text
        left = Inches(5)
        top = Inches(3)
        height = Inches(2)
        width = Inches(3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        p = tf.add_paragraph()
        if len(new_text.split()) > 15:
            third = len(new_text.split()) // 3
            text = ' '.join(new_text.split()[:third] + ["\n"]+ new_text.split()[third:third*2] + ["\n"] + new_text.split()[third*2:])
        else:
            mid = len(new_text.split()) // 2
            text = ' '.join(new_text.split()[:mid] + ["\n"] + new_text.split()[mid:])

        p.text = text
        p.word_wrap = True
        p.font.size = Pt(10)

        # New Image
        left = Inches(0.5)
        top = Inches(5)
        height = Inches(2)
        pic = slide.shapes.add_picture(new_img_path, left, top, height=height)

        # New Image Logits
        left = Inches(0.5)
        top = Inches(4.5)
        height = Inches(1)
        width = Inches(3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
    
        p.text = f"New Image Logits: {dg_lgts:.3f} (Delta: {dg_lgts - orig_lgts:.3f})"
        p.word_wrap = True
        p.font.size = Pt(10)


        # GT Image
        left = Inches(3.75)
        top = Inches(5)
        height = Inches(2)
        pic = slide.shapes.add_picture(gt_img_path, left, top, height=height)

        # GT Image Logits
        left = Inches(3.75)
        top = Inches(4.5)
        height = Inches(1)
        width = Inches(3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()

        p.text = f"Ground Truth Logits: {gt_lgts:.3f} (Delta: {gt_lgts - orig_lgts:.3f})"
        p.word_wrap = True
        p.font.size = Pt(10)

        # Random Image
        left = Inches(7)
        top = Inches(5)
        height = Inches(2)
        pic = slide.shapes.add_picture(random_img_path, left, top, height=height)

        # Random Image Logits
        left = Inches(7)
        top = Inches(4.5)
        height = Inches(1)
        width = Inches(3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()

        p.text = f"Random Image Logits: {rd_lgts:.3f} (Delta: {rd_lgts - orig_lgts:.3f})"
        p.word_wrap = True
        p.font.size = Pt(10)

prs.save('test.pptx')
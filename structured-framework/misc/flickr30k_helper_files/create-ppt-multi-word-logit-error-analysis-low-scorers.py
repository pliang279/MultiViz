from pptx import Presentation, enum
from pptx.util import Inches, Pt
import json

# CLIP Flickr30k Analysis - with Double Grad

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

from structured_framework.misc.flickr30k_clip_target_ids_low_scorers import *
# target_ids_808 = {0: '[CLS]', 1: 'a', 2: 'homeless', 3: 'man', 4: 'is', 5: 'holding', 6: 'a', 7: 'cardboard', 8: 'sign', 9: 'waiting', 10: 'for', 11: 'a', 12: 'ride', 13: 'and', 14: 'possibly', 15: 'money', 16: '.', 17: '[SEP]', 18: '[PAD]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_808 = {
#     "808_1": {"ids": [2], "text": "homeless"},
#     "808_2": {"ids": [3], "text": "man"},
#     "808_3": {"ids": [7], "text": "cardboard"},
#     "808_4": {"ids": [8], "text": "sign"},
#     "808_5": {"ids": [7, 8], "text": "cardboard sign"},
#     "808_6": {"ids": [12], "text": "ride"},
#     "808_7": {"ids": [15], "text": "money"},
#     "808_8": {"ids": [2, 3], "text": "homeless man"},
#     "808_9": {"ids": [2, 3, 4, 5, 6, 7, 8], "text": "homeles man is holding a cardboard sign"}
# }

# # Repeat for:
# # ids = [204,
# # 654,
# # 589,
# # 711,
# # 777,
# # 411,
# # 265,
# # 169,
# # 308,
# # 259,
# # 634,
# # 391,
# # 576]
# target_ids_204 = {0: '[CLS]', 1: 'a', 2: 'indian', 3: 'man', 4: 'with', 5: 'a', 6: 'hairy', 7: 'chest', 8: 'and', 9: 'red', 10: 'eyes', 11: 'is', 12: 'sitting', 13: 'alone', 14: '.', 15: '[SEP]', 16: '[PAD]', 17: '[PAD]', 18: '[PAD]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_204 = {
#     "204_1": {"ids": [2], "text": "indian"},
#     "204_2": {"ids": [3], "text": "man"},
#     "204_3": {"ids": [6], "text": "hairy"},
#     "204_4": {"ids": [7], "text": "chest"},
#     "204_5": {"ids": [9], "text": "red"},
#     "204_6": {"ids": [10], "text": "eyes"},
#     "204_7": {"ids": [3, 4, 5, 6, 7], "text": "man with a hairy chest"},
#     "204_8": {"ids": [3, 4, 5, 6, 7, 8, 9, 10], "text": "man with a hariy chest and red eyes"}
# }

# target_ids_654 = {0: '[CLS]', 1: 'two', 2: 'black', 3: 'and', 4: 'white', 5: 'homeless', 6: 'men', 7: 'smoking', 8: 'and', 9: 'sitting', 10: 'outside', 11: 'next', 12: 'to', 13: 'a', 14: 'building', 15: '.', 16: '[SEP]', 17: '[PAD]', 18: '[PAD]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_654 = {
#     "654_1": {"ids": [2], "text": "black"},
#     "654_2": {"ids": [4], "text": "white"},
#     "654_3": {"ids": [5], "text": "homeless"},
#     "654_4": {"ids": [6], "text": "men"},
#     "654_5": {"ids": [7], "text": "smoking"},
#     "654_6": {"ids": [14], "text": "building"},
#     "654_7": {"ids": [5, 6], "text": "homeless men"},
#     "654_8": {"ids": [5, 6, 7], "text": "homeless men smoking"},
#     "654_9": {"ids": [1, 2, 3, 4, 5, 6], "text": "two black and white homeless men"}
# }

# target_ids_589 = {0: '[CLS]', 1: 'one', 2: 'boy', 3: 'repairing', 4: 'a', 5: 'three', 6: '-', 7: 'wheeled', 8: 'bicycle', 9: 'in', 10: 'a', 11: 'small', 12: 'parking', 13: 'lot', 14: 'while', 15: 'another', 16: 'boy', 17: 'looks', 18: 'on', 19: 'while', 20: 'leaning', 21: 'on', 22: 'the', 23: 'front', 24: 'end', 25: 'of', 26: 'one', 27: 'of', 28: 'the', 29: 'parked', 30: 'cars', 31: '.', 32: '[SEP]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_589 = {
#     "589_1": {"ids": [2], "text": "boy"},
#     "589_2": {"ids": [5, 6, 7], "text": "three-wheeled"},
#     "589_3": {"ids": [8], "text": "bicycle"},
#     "589_4": {"ids": [12], "text": "parking"},
#     "589_5": {"ids": [12, 13], "text": "parking lot"},
#     "589_6": {"ids": [11, 12, 13], "text": "small parking lot"},
#     "589_7": {"ids": [15, 16], "text": "another boy"},
#     "589_8": {"ids": [30], "text": "cars"},
#     "589_9": {"ids": [29, 30], "text": "parked cars"},
#     "589_10": {"ids": [2, 3, 4, 5, 6, 7, 8], "text": "boy repairing a three-wheeled bicycle"},
#     "589_11": {"ids": [5, 6, 7, 8], "text": "three-wheeled bicycle"}
# }

# target_ids_711 = {0: '[CLS]', 1: 'a', 2: 'large', 3: 'man', 4: 'with', 5: 'mp3', 6: 'head', 7: '##phones', 8: 'and', 9: 'a', 10: 'backpack', 11: 'waits', 12: 'at', 13: 'a', 14: 'subway', 15: 'station', 16: 'on', 17: 'wall', 18: 'street', 19: '.', 20: '[SEP]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_711 = {
#     "711_1": {"ids": [2], "text": "large"},
#     "711_2": {"ids": [3], "text": "man"},
#     "711_3": {"ids": [6, 7], "text": "headphones"},
#     "711_4": {"ids": [5, 6, 7], "text": "mp3 headphones"},
#     "711_5": {"ids": [10], "text": "backpack"},
#     "711_6": {"ids": [14], "text": "subway"},
#     "711_7": {"ids": [15], "text": "station"},
#     "711_8": {"ids": [18], "text": "street"},
#     "711_9": {"ids": [17, 18], "text": "wall street"},
#     "711_10": {"ids": [2, 3], "text": "large man"},
#     "711_11": {"ids": [2, 3, 4, 5, 6, 7], "text": "large man with mp3 headphones"},
#     "711_12": {"ids": [18], "text": "street"}
# }

# target_ids_777 = {0: '[CLS]', 1: 'a', 2: 'young', 3: 'woman', 4: 'with', 5: 'bare', 6: 'feet', 7: 'on', 8: 'a', 9: 'crowded', 10: 'road', 11: 'looking', 12: 'down', 13: 'while', 14: 'many', 15: 'people', 16: 'are', 17: 'walking', 18: 'by', 19: 'and', 20: 'some', 21: 'are', 22: 'sleeping', 23: 'and', 24: 'a', 25: 'young', 26: 'man', 27: 'in', 28: 'a', 29: 'black', 30: 'shirt', 31: 'with', 32: 'yellow', 33: 'writing', 34: 'is', 35: 'looking', 36: 'at', 37: 'the', 38: 'people', 39: '[SEP]'}

# instance_text_target_ids_777 = {
#     "777_1": {"ids": [2], "text": "young"},
#     "777_2": {"ids": [3], "text": "woman"},
#     "777_3": {"ids": [6], "text": "feet"},
#     "777_4": {"ids": [5, 6], "text": "bare feet"},
#     "777_5": {"ids": [10], "text": "road"},
#     "777_6": {"ids": [15], "text": "people"},
#     "777_7": {"ids": [26], "text": "man"},
#     "777_8": {"ids": [25, 26], "text": "young man"},
#     "777_9": {"ids": [29, 30], "text": "black shirt"},
#     "777_10": {"ids": [32, 33], "text": "yellow writing"},
#     "777_11": {"ids": [25, 26, 27, 28, 29, 30], "text": "young man in a black shirt"},
#     "777_12": {"ids": [2, 3, 4, 5, 6], "text": "young woman with bare feet"}
# }
# # 411
# # {0: '[CLS]', 1: 'a', 2: 'man', 3: 'and', 4: 'woman', 5: 'sit', 6: 'in', 7: 'a', 8: 'brightly', 9: 'lit', 10: 'stage', 11: 'set', 12: '.', 13: '[SEP]', 14: '[PAD]', 15: '[PAD]', 16: '[PAD]', 17: '[PAD]', 18: '[PAD]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}
# # Key:  411_1
# # [] []
# # 265
# # {0: '[CLS]', 1: 'there', 2: 'are', 3: 'three', 4: 'men', 5: ',', 6: 'one', 7: 'man', 8: 'is', 9: 'pointing', 10: ',', 11: 'the', 12: 'other', 13: 'two', 14: 'men', 15: 'are', 16: 'wearing', 17: 'glasses', 18: 'and', 19: 'a', 20: 'hat', 21: '.', 22: '[SEP]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}
# # Key:  265_1
# # [] []
# # 169
# # {0: '[CLS]', 1: 'a', 2: 'woman', 3: 'carrying', 4: 'another', 5: 'woman', 6: 'wearing', 7: 'matching', 8: 'exercise', 9: 'attire', 10: 'on', 11: 'her', 12: 'back', 13: '.', 14: '[SEP]', 15: '[PAD]', 16: '[PAD]', 17: '[PAD]', 18: '[PAD]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}
# # Key:  169_1
# # [] []
# # 308
# # {0: '[CLS]', 1: 'four', 2: 'young', 3: 'children', 4: 'planning', 5: 'on', 6: 'how', 7: 'they', 8: 'will', 9: 'spend', 10: 'their', 11: 'summer', 12: '.', 13: '[SEP]', 14: '[PAD]', 15: '[PAD]', 16: '[PAD]', 17: '[PAD]', 18: '[PAD]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}
# # Key:  308_1
# # [] []
# # 259
# # {0: '[CLS]', 1: 'a', 2: 'man', 3: 'stretching', 4: 'his', 5: 'arms', 6: 'up', 7: 'at', 8: 'a', 9: 'table', 10: 'with', 11: 'two', 12: 'ladies', 13: 'talking', 14: 'to', 15: 'each', 16: 'other', 17: '.', 18: '[SEP]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}
# # Key:  259_1
# # [] []
# # 634
# # {0: '[CLS]', 1: 'two', 2: 'men', 3: ',', 4: 'likely', 5: 'a', 6: 'father', 7: 'and', 8: 'sun', 9: ',', 10: 'assemble', 11: 'a', 12: 'stainless', 13: 'steel', 14: 'grill', 15: 'together', 16: '.', 17: '[SEP]', 18: '[PAD]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}
# # Key:  634_1
# # [] []
# # 391
# # {0: '[CLS]', 1: 'three', 2: 'dogs', 3: ',', 4: 'one', 5: 'black', 6: 'and', 7: 'two', 8: 'blond', 9: ',', 10: 'run', 11: 'in', 12: 'the', 13: 'snow', 14: '.', 15: '[SEP]', 16: '[PAD]', 17: '[PAD]', 18: '[PAD]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}
# # Key:  391_1
# # [] []
# # 576
# # {0: '[CLS]', 1: 'runners', 2: 'pass', 3: 'a', 4: 'statue', 5: 'of', 6: 'a', 7: 'man', 8: 'on', 9: 'a', 10: 'horse', 11: 'as', 12: 'spectators', 13: 'look', 14: 'on', 15: 'and', 16: 'take', 17: 'photos', 18: '.', 19: '[SEP]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}
# # Key:  576_1

# # target_ids_411 = {
# #     0: "[CLS]",
# #     1: "a",
# #     2: "large"
# # }

# # instance_text_target_ids_411 = {
# #     "411_1": {"ids": [2], "text": "large"}
# # }

# # target_ids_265 = {
# #     0: "[CLS]",
# #     1: "a",
# #     2: "large"
# # }

# # instance_text_target_ids_265 = {
# #     "265_1": {"ids": [2], "text": "large"}
# # }

# # target_ids_169 = {
# #     0: "[CLS]",
# #     1: "a",
# #     2: "large"
# # }

# # instance_text_target_ids_169 = {
# #     "169_1": {"ids": [2], "text": "large"}
# # }

# # target_ids_308 = {
# #     0: "[CLS]",
# #     1: "a",
# #     2: "large"
# # }

# # instance_text_target_ids_308 = {
# #     "308_1": {"ids": [2], "text": "large"}
# # }

# # target_ids_259 = {
# #     0: "[CLS]",
# #     1: "a",
# #     2: "large"
# # }

# # instance_text_target_ids_259 = {
# #     "259_1": {"ids": [2], "text": "large"}
# # }

# # target_ids_634 = {
# #     0: "[CLS]",
# #     1: "a",
# #     2: "large"
# # }

# # instance_text_target_ids_634 = {
# #     "634_1": {"ids": [2], "text": "large"}
# # }

# target_ids_391 = {0: '[CLS]', 1: 'three', 2: 'dogs', 3: ',', 4: 'one', 5: 'black', 6: 'and', 7: 'two', 8: 'blond', 9: ',', 10: 'run', 11: 'in', 12: 'the', 13: 'snow', 14: '.', 15: '[SEP]', 16: '[PAD]', 17: '[PAD]', 18: '[PAD]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_391 = {
#     "391_1": {"ids": [2], "text": "dogs"},
#     "391_2": {"ids": [5], "text": "black"},
#     "391_3": {"ids": [1, 2], "text": "three dogs"}, 
#     "391_4": {"ids": [8], "text": "blonde"},
#     "391_5": {"ids": [7, 8], "text": "two blond"},
#     "391_6": {"ids": [13], "text": "snow"},
#     "391_7": {"ids": [4, 5], "text": "one black"},
#     "391_8": {"ids": [1, 2, 3, 4, 5, 6, 7], "text": "three dogs, one black and two blond"},
# }

# target_ids_576 = {0: '[CLS]', 1: 'runners', 2: 'pass', 3: 'a', 4: 'statue', 5: 'of', 6: 'a', 7: 'man', 8: 'on', 9: 'a', 10: 'horse', 11: 'as', 12: 'spectators', 13: 'look', 14: 'on', 15: 'and', 16: 'take', 17: 'photos', 18: '.', 19: '[SEP]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_576 = {
#     "576_1": {"ids": [1], "text": "runners"},
#     "576_2": {"ids": [4], "text": "statue"},
#     "576_3": {"ids": [7], "text": "man"},
#     "576_4": {"ids": [10], "text": "horse"},
#     "576_5": {"ids": [12], "text": "spectators"},
#     "576_6": {"ids": [4, 5, 6, 7], "text": "statue of a man"},
#     "576_7": {"ids": [4, 5, 6, 7, 8, 9, 10], "text": "statue of a man on a horse"},
#     "576_8": {"ids": [17], "text": "photos"}
# }



id_to_tids = {
    # 50: instance_text_target_ids_50,
    808: instance_text_target_ids_808,
    204: instance_text_target_ids_204,
    654: instance_text_target_ids_654,
    589: instance_text_target_ids_589,
    711: instance_text_target_ids_711,
    777: instance_text_target_ids_777,
    # 411: instance_text_target_ids_411,
    # 265: instance_text_target_ids_265,
    # 169: instance_text_target_ids_169,
    # 308: instance_text_target_ids_308,
    # 259: instance_text_target_ids_259,
    # 634: instance_text_target_ids_634,
    391: instance_text_target_ids_391,
    576: instance_text_target_ids_576
    # 403
}

for instance_idx, tid_dict in id_to_tids.items():
    # # Read Key to Logits JSON
    with open(f"clip_key_to_logits/key_to_logits-box-acc-{instance_idx}.json") as f:
        key_to_logits = json.loads(f.read())
    for key, value in tid_dict.items(): 
        orig_img_path = f"structured_framework/visuals/flickr30k-clip-{instance_idx}-{0}-image.png"
        orig_text = open(f"structured_framework/visuals/flickr30k-clip-{instance_idx}-{0}-text.txt").read()
        doublegrad_img_path = f"structured_framework/visuals/flickr30k-clip-{key}-doublegrad.png"
        new_img_path = f"structured_framework/visuals/flickr30k-clip-{instance_idx}-0-saliency.png"
        # new_text = open(f"structured_framework/visuals/flickr30k-clip-{key}-new_text.txt").read()
        random_img_path =  f"structured_framework/visuals/flickr30k-clip-{key}-new_box_img_with_boxes.jpg"
        # gt_img_path = f"structured_framework/visuals/flickr30k-clip-{key}-gt_img.jpg"
        slide = prs.slides.add_slide(blank_slide_layout)


        # "original_logits": 10.043411,
        # "doublegrad_logits": 3.0287437,
        # "random_drop_logits": 8.995966,
        # "ground_truth_logits": 8.610129

        orig_lgts = key_to_logits[str(instance_idx)][key]["original_logits"]
        # dg_lgts = key_to_logits[str(instance_idx)][key]["doublegrad_box_logits"]
        # rd_lgts = key_to_logits[str(instance_idx)][key]["random_box_logits"]
        # gt_lgts = key_to_logits[str(instance_idx)][key]["ground_truth_logits"]


        # Original Image
        left = top = Inches(0.5)
        height = Inches(2.5)
        pic = slide.shapes.add_picture(orig_img_path, left, top, height=height)
        
        # Original Text
        left = Inches(1)
        top = Inches(3)
        height = Inches(2)
        width = Inches(3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        p = tf.add_paragraph()
        if len(orig_text.split()) > 15:
            third = len(orig_text.split()) // 3
            text = ' '.join(orig_text.split()[:third] + ["\n"]+ orig_text.split()[third:third*2] + ["\n"] + orig_text.split()[third*2:])
        else:
            mid = len(orig_text.split()) // 2
            text = ' '.join(orig_text.split()[:mid] + ["\n"] + orig_text.split()[mid:])
        p.text = text
        p.word_wrap = True
        p.font.size = Pt(10)

        left = Inches(1)
        top = Inches(3.5)
        height = Inches(0.5)
        width = Inches(3)
        txBox2 = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox2.text_frame

        double_grad_text = value["text"]

        if len(double_grad_text.split()) > 15:
            third = len(double_grad_text.split()) // 3
            text = ' '.join(double_grad_text.split()[:third] + ["\n"]+ double_grad_text.split()[third:third*2] + ["\n"] + double_grad_text.split()[third*2:])
        elif len(double_grad_text.split()) > 10:
            mid = len(orig_text.split()) // 2
            text = ' '.join(orig_text.split()[:mid] + ["\n"] + orig_text.split()[mid:])
        else:
            text = double_grad_text

        p = tf.add_paragraph()
        p.text = f"Double Grad Text: {text} \n Orig Logits: {orig_lgts:.3f}"
        p.word_wrap = True
        p.font.size = Pt(10)

        # DoubleGrad Image
        left = Inches(5)
        top = Inches(0.5)
        height = Inches(3)
        pic = slide.shapes.add_picture(doublegrad_img_path, left, top, height=height)

        # # New Text
        # left = Inches(5)
        # top = Inches(3)
        # height = Inches(2)
        # width = Inches(3)
        # txBox = slide.shapes.add_textbox(left, top, width, height)
        # tf = txBox.text_frame

        # p = tf.add_paragraph()
        # if len(new_text.split()) > 15:
        #     third = len(new_text.split()) // 3
        #     text = ' '.join(new_text.split()[:third] + ["\n"]+ new_text.split()[third:third*2] + ["\n"] + new_text.split()[third*2:])
        # else:
        #     mid = len(new_text.split()) // 2
        #     text = ' '.join(new_text.split()[:mid] + ["\n"] + new_text.split()[mid:])

        # p.text = text
        # p.word_wrap = True
        # p.font.size = Pt(10)

        # New Image
        left = Inches(0.5)
        top = Inches(4.5)
        height = Inches(3)
        pic = slide.shapes.add_picture(new_img_path, left, top, height=height)

        # New Image Logits
        left = Inches(0.5)
        top = Inches(4)
        height = Inches(1)
        width = Inches(3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
    
        p.text = f"Unimodal Gradient"
        p.word_wrap = True
        p.font.size = Pt(10)


        # # GT Image
        # left = Inches(3.75)
        # top = Inches(5)
        # height = Inches(2)
        # pic = slide.shapes.add_picture(gt_img_path, left, top, height=height)

        # # GT Image Logits
        # left = Inches(3.75)
        # top = Inches(4.5)
        # height = Inches(1)
        # width = Inches(3)
        # txBox = slide.shapes.add_textbox(left, top, width, height)
        # tf = txBox.text_frame
        # p = tf.add_paragraph()

        # p.text = f"Ground Truth Logits: {gt_lgts:.3f} (Delta: {gt_lgts - orig_lgts:.3f})"
        # p.word_wrap = True
        # p.font.size = Pt(10)

        # Random Image
        left = Inches(7)
        top = Inches(5)
        height = Inches(2)
        pic = slide.shapes.add_picture(random_img_path, left, top, height=height)

        # # Random Image Logits
        # left = Inches(7)
        # top = Inches(4.5)
        # height = Inches(1)
        # width = Inches(3)
        # txBox = slide.shapes.add_textbox(left, top, width, height)
        # tf = txBox.text_frame
        # p = tf.add_paragraph()

        # p.text = f"Random Image Logits: {rd_lgts:.3f} (Delta: {rd_lgts - orig_lgts:.3f})"
        # p.word_wrap = True
        # p.font.size = Pt(10)

prs.save('test.pptx')
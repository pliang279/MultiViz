from pptx import Presentation, enum
from pptx.util import Inches, Pt

# CLIP Flickr30k Analysis - with Double Grad

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

instance_text_target_ids = {
    50: {"ids": [1, 2, 3, 4, 5, 6], "text": "two black and white homeless men"},
    100: {"ids": [1, 2], "text": "the car"},
    150: {"ids": [1, 2], "text": "two dogs"},

    200: {"ids": [13, 14, 15, 16], "text": "shallow wading pool"},
    250: {"ids": [2, 3, 4], "text": "soccer team player"},
    300: {"ids": [1, 2, 3, 4, 5], "text": "two boys, two girls"},
    350: {"ids": [2, 3], "text": "little girl"},
    400: {"ids": [8, 9], "text": "black necklace"},
    450: {
        "ids": [4, 5, 6],
        "text": "the red shirt",
    },
    500: {"ids": [15, 16], "text": "the foothills"}
}

logits_and_props = {50: {'logits': -9.155246, 'probs': 1.0}, 100: {'logits': -10.669653, 'probs': 1.0}, 150: {'logits': -10.611729, 'probs': 1.0}, 200: {'logits': -10.681751, 'probs': 1.0}, 250: {'logits': -10.703483, 'probs': 1.0}, 300: {'logits': -10.320086, 'probs': 1.0}, 350: {'logits': -10.690347, 'probs': 1.0}, 400: {'logits': -10.634938, 'probs': 1.0}, 450: {'logits': -10.706868, 'probs': 1.0}, 500: {'logits': -10.660023, 'probs': 1.0}}
for instance_idx in [50, 100, 150, 200, 250, 300, 350, 400, 450, 500]:
    orig_img_path = f"structured-framework/visuals/flickr30k-vilt-{instance_idx}-{0}-image.png"
    orig_text = open(f"structured-framework/visuals/flickr30k-vilt-{instance_idx}-{0}-text.txt").read()
    lime_img_path = f"structured-framework/visuals/flickr30k-vilt-{instance_idx}-{0}-image-lime-pred.png"
    saliency_img_path = f"structured-framework/visuals/flickr30k-vilt-{instance_idx}-{0}-saliency.png"
    # lime_text_path = f"structured-framework/visuals/flickr30k-vilt-{instance_idx}-{0}-text-lime-pred.png"
    doublegrad_img_path = f"structured-framework/visuals/flickr30k-vilt-{instance_idx}-{0}-doublegrad.png"
    slide = prs.slides.add_slide(blank_slide_layout)

    # Original Image
    left = top = Inches(0.5)
    height = Inches(2.5)
    pic = slide.shapes.add_picture(orig_img_path, left, top, height=height)
    
    # Original Text
    left = Inches(1)
    top = Inches(3)
    height = Inches(2)
    width = Inches(3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.add_paragraph()
    if len(orig_text.split()) > 15:
        third = len(orig_text.split()) // 3
        text = ' '.join(orig_text.split()[:third] + ["\n"]+ orig_text.split()[third:third*2] + ["\n"] + orig_text.split()[third*2:])
    else:
        mid = len(orig_text.split()) // 2
        text = ' '.join(orig_text.split()[:mid] + ["\n"] + orig_text.split()[mid:])
    p.text = text
    p.word_wrap = True
    p.font.size = Pt(10)

    left = Inches(1)
    top = Inches(3.5)
    height = Inches(0.5)
    width = Inches(3)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox2.text_frame

    double_grad_text = instance_text_target_ids[instance_idx]["text"] + f' Logits: {logits_and_props[instance_idx]["logits"]:.2f}, Probs: {logits_and_props[instance_idx]["probs"]}'

    if len(double_grad_text.split()) > 15:
        third = len(double_grad_text.split()) // 3
        text = ' '.join(double_grad_text.split()[:third] + ["\n"]+ double_grad_text.split()[third:third*2] + ["\n"] + double_grad_text.split()[third*2:])
    elif len(double_grad_text.split()) > 10:
        mid = len(orig_text.split()) // 2
        text = ' '.join(orig_text.split()[:mid] + ["\n"] + orig_text.split()[mid:])
    else:
        text = double_grad_text

    p = tf.add_paragraph()
    p.text = "Double Grad Text: " + text
    p.word_wrap = True
    p.font.size = Pt(10)


    # LIME Text
    left = Inches(0.5)
    top = Inches(4.5)
    height = Inches(3)
    pic = slide.shapes.add_picture(doublegrad_img_path, left, top, height=height)


    # LIME Image
    left = Inches(5)
    top = Inches(0)
    height = Inches(3.5)
    pic = slide.shapes.add_picture(lime_img_path, left, top, height=height)

    # Saliency Image
    left = Inches(5)
    top = Inches(3.5)
    height = Inches(3.5)
    pic = slide.shapes.add_picture(saliency_img_path, left, top, height=height)

    prs.save('test.pptx')
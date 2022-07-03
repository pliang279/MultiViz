from pptx import Presentation, enum
from pptx.util import Inches, Pt
import json

# CLIP Flickr30k Analysis - with Double Grad

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

from structured_framework.misc.flickr30k_clip_target_ids import *
# target_ids_100 = {
#     0: "[CLS]",
#     1: "a",
#     2: "large",
#     3: "bearded",
#     4: "man",
#     5: "flip",
#     6: "##s",
#     7: "a",
#     8: "cr",
#     9: "##ep",
#     10: "##e",
#     11: "or",
#     12: "om",
#     13: "##ele",
#     14: "##t",
#     15: "in",
#     16: "mid",
#     17: "##air",
#     18: "with",
#     19: "his",
#     20: "fry",
#     21: "##ing",
#     22: "pan",
#     23: ".",
#     24: "[SEP]",
#     25: "[PAD]",
#     26: "[PAD]",
#     27: "[PAD]",
#     28: "[PAD]",
#     29: "[PAD]",
#     30: "[PAD]",
#     31: "[PAD]",
#     32: "[PAD]",
#     33: "[PAD]",
#     34: "[PAD]",
#     35: "[PAD]",
#     36: "[PAD]",
#     37: "[PAD]",
#     38: "[PAD]",
#     39: "[PAD]",
# }

# instance_text_target_ids_100 = {
#     "100_1": {"ids": [2], "text": "large"},
#     "100_2": {"ids": [3], "text": "bearded"},
#     "100_3": {"ids": [4], "text": "man"},
#     "100_4": {"ids": [8, 9, 10], "text": "crepe"},
#     "100_5": {"ids": [12, 13, 14], "text": "omelet"},
#     "100_6": {"ids": [20, 21, 22], "text": "frying pan"},
#     "100_7": {"ids": [2, 3, 4], "text": "large bearded man"},
#     "100_8": {
#         "ids": [2, 3, 4, 5, 6, 7, 8, 9, 10],
#         "text": "large bearded man flips a crepe",
#     },
#     "100_9": {
#         "ids": [2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12, 13, 14],
#         "text": "large bearded man flips a crepe or omelet",
#     },
#     "100_10": {"ids": [12, 13, 14, 15, 16, 17], "text": "omelet in mid air"},
# }

# target_ids_150 = {
#     0: "[CLS]",
#     1: "a",
#     2: "black",
#     3: "dog",
#     4: "with",
#     5: "white",
#     6: "facial",
#     7: "and",
#     8: "chest",
#     9: "markings",
#     10: "standing",
#     11: "in",
#     12: "chest",
#     13: "high",
#     14: "water",
#     15: ".",
#     16: "[SEP]",
#     17: "[PAD]",
#     18: "[PAD]",
#     19: "[PAD]",
#     20: "[PAD]",
#     21: "[PAD]",
#     22: "[PAD]",
#     23: "[PAD]",
#     24: "[PAD]",
#     25: "[PAD]",
#     26: "[PAD]",
#     27: "[PAD]",
#     28: "[PAD]",
#     29: "[PAD]",
#     30: "[PAD]",
#     31: "[PAD]",
#     32: "[PAD]",
#     33: "[PAD]",
#     34: "[PAD]",
#     35: "[PAD]",
#     36: "[PAD]",
#     37: "[PAD]",
#     38: "[PAD]",
#     39: "[PAD]",
# }

# instance_text_target_ids_150 = {
#     "150_1": {"ids": [2], "text": "black"},
#     "150_2": {"ids": [3], "text": "dog"},
#     "150_3": {"ids": [5], "text": "white"},
#     "150_4": {"ids": [6], "text": "facial"},
#     "150_5": {"ids": [8], "text": "chest"},
#     "150_6": {"ids": [9], "text": "markings"},
#     "150_7": {"ids": [10], "text": "standing"},
#     "150_8": {"ids": [11], "text": "in"},
#     "150_9": {"ids": [12], "text": "chest"},
#     "150_10": {"ids": [13], "text": "high"},
#     "150_11": {"ids": [14], "text": "water"},
#     "150_12": {"ids": [2, 3], "text": "black dog"},
#     "150_13": {"ids": [5, 6], "text": "white facial"},
#     "150_14": {"ids": [5, 6, 7, 8, 9], "text": "white facial and chest markings"},
#     "150_15": {"ids": [12, 13, 14], "text": "chest high water"},
# }

# target_ids_200 = {
#     0: "[CLS]",
#     1: "a",
#     2: "man",
#     3: "is",
#     4: "taking",
#     5: "photographs",
#     6: "of",
#     7: "a",
#     8: "large",
#     9: "garden",
#     10: "of",
#     11: "white",
#     12: "and",
#     13: "orange",
#     14: "tu",
#     15: "##lip",
#     16: "##s",
#     17: ".",
#     18: "[SEP]",
#     19: "[PAD]",
#     20: "[PAD]",
#     21: "[PAD]",
#     22: "[PAD]",
#     23: "[PAD]",
#     24: "[PAD]",
#     25: "[PAD]",
#     26: "[PAD]",
#     27: "[PAD]",
#     28: "[PAD]",
#     29: "[PAD]",
#     30: "[PAD]",
#     31: "[PAD]",
#     32: "[PAD]",
#     33: "[PAD]",
#     34: "[PAD]",
#     35: "[PAD]",
#     36: "[PAD]",
#     37: "[PAD]",
#     38: "[PAD]",
#     39: "[PAD]",
# }

# instance_text_target_ids_200 = {
#     "200_1": {"ids": [2], "text": "man"},
#     "200_2": {"ids": [5], "text": "photographs"},
#     "200_3": {"ids": [9], "text": "garden"},
#     "200_4": {"ids": [11], "text": "white"},
#     "200_5": {"ids": [13], "text": "orange"},
#     "200_6": {"ids": [14, 15, 16], "text": "tulips"},
#     "200_7": {"ids": [1, 2, 3, 4, 5], "text": "a man is taking photographs"},
#     "200_8": {"ids": [8, 9], "text": "large garden"},
#     "200_9": {
#         "ids": [1, 2, 3, 4, 5, 6, 7, 8, 9],
#         "text": "a man is taking photographs of a large garden",
#     },
#     "200_10": {
#         "ids": [8, 9, 10, 11, 12, 13, 14, 15, 16],
#         "text": "a large garden of white and orange tulips",
#     },
#     "200_11": {"ids": [11, 12, 13, 14, 15, 16], "text": "white and orange tulips"},
#     "200_12": {
#         "ids": [1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12, 13, 14, 15, 16],
#         "text": "a man is taking photographs of a large garden of white and orange tulips",
#     },
# }

# target_ids_500 = {
#     0: "[CLS]",
#     1: "a",
#     2: "little",
#     3: "girl",
#     4: "in",
#     5: "front",
#     6: "a",
#     7: "pink",
#     8: "food",
#     9: "tray",
#     10: "is",
#     11: "getting",
#     12: "her",
#     13: "bike",
#     14: "helmet",
#     15: "on",
#     16: "by",
#     17: "a",
#     18: "woman",
#     19: ".",
#     20: "[SEP]",
#     21: "[PAD]",
#     22: "[PAD]",
#     23: "[PAD]",
#     24: "[PAD]",
#     25: "[PAD]",
#     26: "[PAD]",
#     27: "[PAD]",
#     28: "[PAD]",
#     29: "[PAD]",
#     30: "[PAD]",
#     31: "[PAD]",
#     32: "[PAD]",
#     33: "[PAD]",
#     34: "[PAD]",
#     35: "[PAD]",
#     36: "[PAD]",
#     37: "[PAD]",
#     38: "[PAD]",
#     39: "[PAD]",
# }

# instance_text_target_ids_500 = {
#     "500_1": {"ids": [2], "text": "little"},
#     "500_2": {"ids": [3], "text": "girl"},
#     "500_3": {"ids": [7], "text": "pink"},
#     "500_4": {"ids": [8], "text": "food"},
#     "500_5": {"ids": [9], "text": "tray"},
#     "500_6": {"ids": [13], "text": "bike"},
#     "500_7": {"ids": [14], "text": "helmet"},
#     "500_8": {"ids": [18], "text": "woman"},
#     "500_9": {"ids": [2, 3], "text": "little girl"},
#     "500_10": {"ids": [8, 9], "text": "food tray"},
#     "500_11": {"ids": [7, 8, 9], "text": "pink food tray"},
#     "500_12": {
#         "ids": [2, 3, 4, 5, 6, 7, 8, 9],
#         "text": "little girl in front a pink food tray",
#     },
#     "500_13": {"ids": [13, 14], "text": "bike helmet"},
# }

# target_ids_50 = {
#     0: "[CLS]",
#     1: "three",
#     2: "small",
#     3: "dogs",
#     4: ",",
#     5: "two",
#     6: "white",
#     7: "and",
#     8: "one",
#     9: "black",
#     10: "and",
#     11: "white",
#     12: ",",
#     13: "on",
#     14: "a",
#     15: "sidewalk",
#     16: ".",
#     17: "[SEP]",
#     18: "[PAD]",
#     19: "[PAD]",
#     20: "[PAD]",
#     21: "[PAD]",
#     22: "[PAD]",
#     23: "[PAD]",
#     24: "[PAD]",
#     25: "[PAD]",
#     26: "[PAD]",
#     27: "[PAD]",
#     28: "[PAD]",
#     29: "[PAD]",
#     30: "[PAD]",
#     31: "[PAD]",
#     32: "[PAD]",
#     33: "[PAD]",
#     34: "[PAD]",
#     35: "[PAD]",
#     36: "[PAD]",
#     37: "[PAD]",
#     38: "[PAD]",
#     39: "[PAD]",
# }

# instance_text_target_ids_50 = {
#     "50_1": {"ids": [2], "text": "small"},
#     "50_2": {"ids": [3], "text": "dogs"},
#     "50_3": {"ids": [1], "text": "three"},
#     "50_4": {"ids": [2, 3], "text": "small dogs"},
#     "50_5": {"ids": [1, 2, 3], "text": "three small dogs"},
#     "50_6": {"ids": [6], "text": "white"},
#     "50_7": {"ids": [9], "text": "black"},
#     "50_8": {"ids": [9, 10, 11], "text": "black and white"},
#     "50_9": {"ids": [15], "text": "sidewalk"},
#     "50_10": {
#         "ids": [5, 6, 7, 8, 9, 10, 11],
#         "text": "two white and one black and white",
#     },
#     "50_11": {
#         "ids": [1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11],
#         "text": "three small dogs, two white and one black and white",
#     },
# }

# target_ids_250 = {
#     0: "[CLS]",
#     1: "two",
#     2: "boys",
#     3: ",",
#     4: "two",
#     5: "girls",
#     6: ",",
#     7: "strapped",
#     8: "in",
#     9: "and",
#     10: "ready",
#     11: "for",
#     12: "an",
#     13: "amusement",
#     14: "park",
#     15: "ride",
#     16: ".",
#     17: "[SEP]",
#     18: "[PAD]",
#     19: "[PAD]",
#     20: "[PAD]",
#     21: "[PAD]",
#     22: "[PAD]",
#     23: "[PAD]",
#     24: "[PAD]",
#     25: "[PAD]",
#     26: "[PAD]",
#     27: "[PAD]",
#     28: "[PAD]",
#     29: "[PAD]",
#     30: "[PAD]",
#     31: "[PAD]",
#     32: "[PAD]",
#     33: "[PAD]",
#     34: "[PAD]",
#     35: "[PAD]",
#     36: "[PAD]",
#     37: "[PAD]",
#     38: "[PAD]",
#     39: "[PAD]",
# }
# instance_text_target_ids_250 = {
#     "250_1": {"ids": [1], "text": "two"},
#     "250_2": {"ids": [2], "text": "boys"},
#     "250_3": {"ids": [5], "text": "girls"},
#     "250_4": {"ids": [7], "text": "strapped"},
#     "250_5": {"ids": [13], "text": "amusement"},
#     "250_6": {"ids": [14], "text": "park"},
#     "250_7": {"ids": [15], "text": "ride"},
#     "250_8": {"ids": [1, 2], "text": "two boys"},
#     "250_9": {"ids": [4, 5], "text": "two girls"},
#     "250_10": {"ids": [1, 2, 3, 4, 5], "text": "two boys, two girls"},
#     "250_11": {"ids": [13, 14], "text": "amusement park"},
#     "250_12": {"ids": [14, 15], "text": "park ride"},
#     "250_13": {"ids": [13, 14, 15], "text": "amusement park ride"}
# }

# target_ids_300 = {0: '[CLS]', 1: 'a', 2: 'young', 3: 'boy', 4: 'wearing', 5: 'a', 6: 'black', 7: 'shirt', 8: 'and', 9: 'brown', 10: 'pants', 11: 'practices', 12: 'jumping', 13: 'over', 14: 'a', 15: 'low', 16: 'bar', 17: 'on', 18: 'a', 19: 'skate', 20: '##board', 21: '.', 22: '[SEP]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_300 = {
#     "300_1": {"ids": [2], "text": "young"},
#     "300_2": {"ids": [3], "text": "boy"},
#     "300_3": {"ids": [6], "text": "black"},
#     "300_4": {"ids": [7], "text": "shirt"},
#     "300_5": {"ids": [9], "text": "brown"},
#     "300_6": {"ids": [10], "text": "pants"},
#     "300_7": {"ids": [12], "text": "jumping"},
#     "300_8": {"ids": [16], "text": "bar"},
#     "300_9": {"ids": [15, 16], "text": "low bar"},
#     "300_10": {"ids": [19, 20], "text": "skateboard"},
#     "300_11": {"ids": [2, 3], "text": "young boy"},
#     "300_12": {"ids": [6, 7], "text": "black shirt"},
#     "300_13": {"ids": [9, 10], "text": "brown pants"},
#     "300_14": {"ids": [6,7,8,9,10], "text": "black shirt and brown pants"},
#     "300_15": {"ids": [2,3,4,5,6,7,8,9,10], "text": "young boy wearing a black shirt and brown pants"}
# }

# target_ids_400 = {0: '[CLS]', 1: 'a', 2: 'woman', 3: 'in', 4: 'a', 5: 'jean', 6: 'jacket', 7: 'and', 8: 'black', 9: 'sunglasses', 10: 'stands', 11: 'outside', 12: 'with', 13: 'two', 14: 'young', 15: 'boys', 16: 'by', 17: 'a', 18: 'ki', 19: '##os', 20: '##k', 21: ',', 22: 'looking', 23: 'at', 24: 'a', 25: 'paper', 26: 'she', 27: 'is', 28: 'holding', 29: 'in', 30: 'her', 31: 'hand', 32: '.', 33: '[SEP]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_400 = {
#     "400_1": {"ids": [2], "text": "woman"},
#     "400_2": {"ids": [5], "text": "jean"},
#     "400_3": {"ids": [6], "text": "jacket"},
#     "400_4": {"ids": [8], "text": "black"},
#     "400_5": {"ids": [9], "text": "sunglasses"},
#     "400_6": {"ids": [13], "text": "two"},
#     "400_7": {"ids": [14], "text": "young"},
#     "400_8": {"ids": [15], "text": "boys"},
#     "400_9": {"ids": [18, 19, 20], "text": "kiosk"},
#     "400_10": {"ids": [5, 6], "text": "jean jacket"},
#     "400_11": {"ids": [8, 9], "text": "black sunglasses"},
#     "400_12": {"ids": [13, 14, 15], "text": "two young boys"},
#     "400_13": {"ids": [2, 3, 4, 5, 6], "text": "woman in a jean jacket"},
#     "400_14": {"ids": [2, 3, 4, 5, 6, 7, 8, 9], "text": "woman in a jean jacket and black sunglasses"},
#     "400_15": {"ids": [25], "text": "paper"},
#     "400_16": {"ids": [31], "text": "hand"}
# }

# target_ids_600 = {0: '[CLS]', 1: 'a', 2: 'hooded', 3: 'individual', 4: 'with', 5: 'an', 6: 'orange', 7: 'scarf', 8: 'and', 9: 'face', 10: 'covering', 11: 'uses', 12: 'a', 13: 'small', 14: 'knife', 15: 'to', 16: 'sc', 17: '##ul', 18: '##pt', 19: 'a', 20: 'piece', 21: 'of', 22: 'ice', 23: '.', 24: '[SEP]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_600 = {
#     "600_1": {"ids": [2], "text": "hooded"},
#     "600_2": {"ids": [3], "text": "individual"},
#     "600_3": {"ids": [2, 3], "text": "hooden individual"},
#     "600_4": {"ids": [6], "text": "orange"},
#     "600_5": {"ids": [7], "text": "scarf"},
#     "600_6": {"ids": [9], "text": "face"},
#     "600_7": {"ids": [10], "text": "covering"},
#     "600_8": {"ids": [14], "text": "knife"},
#     "600_9": {"ids": [16, 17, 18], "text": "sculpt"},
#     "600_10": {"ids": [20], "text": "piece"},
#     "600_11": {"ids": [22], "text": "ice"},
#     "600_12": {"ids": [6, 7], "text": "orange scarf"},
#     "600_13": {"ids": [9, 10], "text": "face covering"},
#     "600_14": {"ids": [20, 21, 22], "text": "piece of ice"},
#     "600_15": {"ids": [2, 3, 4, 5, 6, 7], "text": "hooden individual with an orange scarf"},
#     "600_16": {"ids": [2, 3, 4, 5, 6, 7, 8, 9, 10], "text": "hooden individual with an orange scarf and face covering"},
#     "600_17": {"ids": [13, 14], "text": "small knife"},
# }

# target_ids_700 = {0: '[CLS]', 1: 'guy', 2: 'in', 3: 'jeans', 4: 'and', 5: 'black', 6: 'jacket', 7: 'walking', 8: 'along', 9: 'grass', 10: 'and', 11: 'trees', 12: 'with', 13: 'the', 14: 'city', 15: 'in', 16: 'the', 17: 'background', 18: '.', 19: '[SEP]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_700 = {
#     "700_1": {"ids": [1], "text": "guy"},
#     "700_2": {"ids": [3], "text": "jeans"},
#     "700_3": {"ids": [5], "text": "black"},
#     "700_4": {"ids": [6], "text": "jacket"},
#     "700_5": {"ids": [9], "text": "grass"},
#     "700_6": {"ids": [11], "text": "trees"},
#     "700_7": {"ids": [14], "text": "city"},
#     "700_8": {"ids": [17], "text": "background"},
#     "700_9": {"ids": [1, 2, 3], "text": "guy in jeans"},
#     "700_10": {"ids": [5, 6], "text": "black jacket"},
#     "700_11": {"ids": [9, 10, 11], "text": "grass and trees"},
#     "700_12": {"ids": [14, 15, 16, 17], "text": "city in the background"},
#     "700_13": {"ids": [1, 2, 3, 4, 5, 6], "text": "guy in jeans and black jacket"}
# }

# target_ids_800 = {0: '[CLS]', 1: 'five', 2: 'children', 3: 'spin', 4: 'around', 5: 'on', 6: 'a', 7: 'playground', 8: 'roundabout', 9: ';', 10: 'three', 11: 'lay', 12: 'on', 13: 'their', 14: 'backs', 15: ',', 16: 'while', 17: 'one', 18: 'attempts', 19: 'to', 20: 'pull', 21: 'himself', 22: 'up', 23: 'with', 24: 'both', 25: 'arms', 26: ',', 27: 'and', 28: 'another', 29: 'holds', 30: 'onto', 31: 'the', 32: 'side', 33: 'while', 34: 'sitting', 35: 'up', 36: '.', 37: '[SEP]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_800 = {
#     "800_1": {"ids": [1], "text": "five"},
#     "800_2": {"ids": [2], "text": "children"},
#     "800_3": {"ids": [7], "text": "playground"},
#     "800_4": {"ids": [8], "text": "roundabout"},
#     "800_5": {"ids": [10], "text": "three"},
#     "800_6": {"ids": [14], "text": "backs"},
#     "800_7": {"ids": [25], "text": "arms"},
#     "800_8": {"ids": [32], "text": "side"},
#     "800_9": {"ids": [34], "text": "sitting"},
#     "800_10": {"ids": [1, 2], "text": "five children"},
#     "800_11": {"ids": [7, 8], "text": "playground roundabout"},
#     "800_12": {"ids": [10, 11, 12, 13, 14], "text": "three lay on their backs"},
#     "800_13": {"ids": [1, 2, 3, 4, 5, 6, 7, 8], "text": "five children spin around on a playground roundabout"},
#     "800_14": {"ids": [17, 18, 19, 20, 21, 22, 23, 24, 25], "text": "one attempts to pull himself up with both arms"},
#     "800_15": {"ids": [28, 29, 30, 31, 32, 33, 34], "text": "another holds onto the side while sitting up"}
# }

# target_ids_900 = {0: '[CLS]', 1: 'a', 2: 'football', 3: 'player', 4: 'preparing', 5: 'a', 6: 'football', 7: 'for', 8: 'a', 9: 'field', 10: 'goal', 11: 'kick', 12: ',', 13: 'while', 14: 'his', 15: 'teammates', 16: 'can', 17: 'coach', 18: 'watch', 19: 'him', 20: '.', 21: '[SEP]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_900 = {
#     "900_1": {"ids": [2], "text": "football"},
#     "900_2": {"ids": [3], "text": "player"},
#     "900_3": {"ids": [6], "text": "football"},
#     "900_4": {"ids": [9, 10], "text": "field goal"},
#     "900_5": {"ids": [11], "text": "kick"},
#     "900_6": {"ids": [9, 10, 11], "text": "field goal kick"},
#     "900_6": {"ids": [15], "text": "teammates"},
#     "900_7": {"ids": [17], "text": "coach"},
#     "900_8": {"ids": [19], "text": "him"},
#     "900_9": {"ids": [15, 16, 17, 18, 19], "text": "teammates can coach watch him"},
#     "900_10": {"ids": [1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11], "text": "a football player preparing a football for a field goal kick"},
# }

# target_ids_1000 = {0: '[CLS]', 1: 'a', 2: 'group', 3: 'of', 4: 'woman', 5: 'from', 6: 'various', 7: 'ethnic', 8: 'backgrounds', 9: 'are', 10: 'competing', 11: 'in', 12: 'a', 13: 'marathon', 14: '.', 15: '[SEP]', 16: '[PAD]', 17: '[PAD]', 18: '[PAD]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_1000 = {
#     "1000_1": {"ids": [2], "text": "group"},
#     "1000_2": {"ids": [4], "text": "woman"},
#     "1000_3": {"ids": [13], "text": "marathon"},
#     "1000_4": {"ids": [6, 7, 8], "text": "various ethnic backgrounds"},
#     "1000_5": {"ids": [2, 3, 4], "text": "group of woman"},
#     "1000_6": {"ids": [2, 3, 4, 5, 6, 7, 8], "text": "group of woman from various ethnic backgrounds"}
# }

# target_ids_350 = {0: '[CLS]', 1: 'a', 2: 'man', 3: 'standing', 4: 'on', 5: 'a', 6: 'street', 7: 'with', 8: 'a', 9: 'suitcase', 10: 'in', 11: 'front', 12: 'of', 13: 'him', 14: 'while', 15: 'another', 16: 'man', 17: 'bends', 18: 'down', 19: 'to', 20: 'look', 21: 'at', 22: 'what', 23: 'is', 24: 'displayed', 25: 'on', 26: 'top', 27: 'of', 28: 'it', 29: '.', 30: '[SEP]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_350 = {
#     "350_1": {"ids": [2], "text": "man"},
#     "350_2": {"ids": [6], "text": "street"},
#     "350_3": {"ids": [9], "text": "suitcase"},
#     "350_4": {"ids": [16], "text": "man"},
#     "350_5": {"ids": [15, 16], "text": "another man"},
#     "350_6": {"ids": [22], "text": "what"},
#     "350_7": {"ids": [2,3,4,5, 6], "text": "man standing on a street"},
#     "350_8": {"ids": [2, 3, 4, 5, 6, 7, 8, 9], "text": "man standing on a street with a suitcase"},
#     "350_9": {"ids": [16, 17, 18], "text": "man bends down"}
# }

# target_ids_450 = {0: '[CLS]', 1: 'a', 2: 'white', 3: 'dog', 4: 'with', 5: 'brown', 6: 'ears', 7: 'is', 8: 'running', 9: 'on', 10: 'the', 11: 'sidewalk', 12: '.', 13: '[SEP]', 14: '[PAD]', 15: '[PAD]', 16: '[PAD]', 17: '[PAD]', 18: '[PAD]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_450 = {
#     "450_1": {"ids": [2], "text": "white"},
#     "450_2": {"ids": [3], "text": "dog"},
#     "450_3": {"ids": [5], "text": "brown"},
#     "450_4": {"ids": [6], "text": "ears"},
#     "450_5": {"ids": [8], "text": "running"},
#     "450_6": {"ids": [11], "text": "sidewalk"},
#     "450_7": {"ids": [2, 3, 4, 5, 6], "text": "white dog with brown ears"},
#     "450_8": {"ids": [2, 3], "text": "white dog"},
#     "450_9": {"ids": [5, 6], "text": "brown ears"},
# }

# target_ids_550 = {0: '[CLS]', 1: 'women', 2: 'sit', 3: 'on', 4: 'a', 5: 'beach', 6: 'as', 7: 'they', 8: 'watch', 9: 'men', 10: 'bring', 11: 'in', 12: 'a', 13: 'fishing', 14: 'net', 15: 'filled', 16: 'with', 17: 'fish', 18: 'in', 19: 'front', 20: 'of', 21: 'them', 22: '.', 23: '[SEP]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_550 = {
#     "550_1": {"ids": [1], "text": "women"},
#     "550_2": {"ids": [5], "text": "beach"},
#     "550_3": {"ids": [9], "text": "men"},
#     "550_4": {"ids": [14], "text": "net"},
#     "550_5": {"ids": [13, 14], "text": "fishing net"},
#     "550_6": {"ids": [17], "text": "fish"},
#     "550_7": {"ids": [1, 2, 3, 4, 5], "text": "women sit on a beach"},
#     "550_8": {"ids": [13, 14, 15, 16, 17], "text": "fishing net filled with fish"}
# }

# target_ids_650 = {0: '[CLS]', 1: 'a', 2: 'man', 3: 'and', 4: 'woman', 5: 'standing', 6: 'against', 7: 'a', 8: 'marble', 9: 'building', 10: ',', 11: 'holding', 12: 'a', 13: 'conversation', 14: 'at', 15: 'night', 16: '.', 17: '[SEP]', 18: '[PAD]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_650 = {
#     "650_1": {"ids": [2], "text": "man"},
#     "650_2": {"ids": [4], "text": "woman"},
#     "650_3": {"ids": [8], "text": "marble"},
#     "650_4": {"ids": [9], "text": "building"},
#     "650_5": {"ids": [15], "text": "night"},
#     "650_6": {"ids": [2, 3, 4], "text": "man and woman"},
#     "650_7": {"ids": [8, 9], "text": "marble building"},
#     "650_8": {"ids": [2, 3, 4, 5, 6, 7, 8, 9], "text": "man and woman standing against a marble building"},
# }

# target_ids_750 = {0: '[CLS]', 1: 'a', 2: 'young', 3: 'man', 4: 'in', 5: 'white', 6: 't', 7: '-', 8: 'shirt', 9: 'is', 10: 'sitting', 11: 'on', 12: 'the', 13: 'floor', 14: 'of', 15: 'a', 16: 'living', 17: 'room', 18: 'full', 19: 'of', 20: 'luggage', 21: '.', 22: '[SEP]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_750 = {
#     "750_1": {"ids": [2], "text": "young"},
#     "750_2": {"ids": [3], "text": "man"},
#     "750_3": {"ids": [5], "text": "white"},
#     "750_4": {"ids": [6, 7, 8], "text": "t-shirt"},
#     "750_5": {"ids": [13], "text": "floor"},
#     "750_6": {"ids": [16, 17], "text": "living room"},
#     "750_7": {"ids": [20], "text": "luggage"},
#     "750_8": {"ids": [2, 3], "text": "young man"},
#     "750_9": {"ids": [2, 3, 4, 5, 7, 8], "text": "young man in white t-shirt"},
#     "750_10": {"ids": [16, 17, 18, 19, 20],"text": "living room full of luggage"},

# }

# target_ids_850 = {0: '[CLS]', 1: 'a', 2: 'baby', 3: ',', 4: 'wearing', 5: 'a', 6: 'pink', 7: 'knit', 8: '##ted', 9: 'hat', 10: ',', 11: 'sleeps', 12: 'peacefully', 13: 'on', 14: 'a', 15: 'blanket', 16: '.', 17: '[SEP]', 18: '[PAD]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_850 = {
#     "850_1": {"ids": [2], "text": "baby"},
#     "850_2": {"ids": [5], "text": "pink"},
#     "850_3": {"ids": [7, 8, 9], "text": "knitted hat"},
#     "850_4": {"ids": [9], "text": "hat"},
#     "850_5": {"ids": [15], "text": "blanket"},
#     "850_6": {"ids": [2, 3, 4, 5, 6, 7, 8, 9], "text": "baby, wearing a pink knitted hat"}
# }

# target_ids_950 = {0: '[CLS]', 1: 'two', 2: 'young', 3: 'girls', 4: 'wearing', 5: 'hi', 6: '##ja', 7: '##bs', 8: 'stand', 9: 'in', 10: 'a', 11: 'dirt', 12: 'courtyard', 13: ',', 14: 'one', 15: 'has', 16: 'her', 17: 'arms', 18: 'folded', 19: 'and', 20: 'looks', 21: 'away', 22: 'the', 23: 'other', 24: 'is', 25: 'staring', 26: 'into', 27: 'the', 28: 'camera', 29: 'with', 30: 'her', 31: 'hands', 32: 'on', 33: 'her', 34: 'hips', 35: '.', 36: '[SEP]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

# instance_text_target_ids_950 = {
#     "950_1": {"ids": [1], "text": "two"},
#     "950_2": {"ids": [2], "text": "young"},
#     "950_3": {"ids": [3], "text": "girls"}, 
#     "950_4": {"ids": [5, 6, 7], "text": "hijabs"},
#     "950_5": {"ids": [12], "text": "courtyard"},
#     "950_6": {"ids": [17], "text": "arms"},
#     "950_7": {"ids": [17, 18], "text": "arms folded"},
#     "950_8": {"ids": [17, 18, 19, 20, 21], "text": "arms folded and looks away"},
#     "950_9": {"ids": [28], "text": "camera"},
#     "950_10": {"ids": [30, 31, 32, 33, 34], "text": "her hands on her hips"},
#     "950_11": {"ids": [31], "text": "hands"},
#     "950_12": {"ids": [34], "text": "hips"},
#     "950_13": {"ids": [30, 31], "text": "her hands"},
#     "950_14": {"ids": [23, 24, 25, 26, 27, 28], "text": "other is staring into the camera"}
# }

id_to_tids = {
    50: instance_text_target_ids_50,
    100: instance_text_target_ids_100,
    150: instance_text_target_ids_150,
    200: instance_text_target_ids_200,
    500: instance_text_target_ids_500,
    250: instance_text_target_ids_250,
    300: instance_text_target_ids_300,
    400: instance_text_target_ids_400,
    600: instance_text_target_ids_600,
    700: instance_text_target_ids_700,
    800: instance_text_target_ids_800,
    900: instance_text_target_ids_900,
    1000: instance_text_target_ids_1000,
    350: instance_text_target_ids_350,
    450: instance_text_target_ids_450,
    550: instance_text_target_ids_550,
    650: instance_text_target_ids_650,
    750: instance_text_target_ids_750,
    850: instance_text_target_ids_850,
    950: instance_text_target_ids_950
}



for instance_idx, tid_dict in id_to_tids.items():
    # # Read Key to Logits JSON
    with open(f"clip_key_to_logits/key_to_logits-box-acc-{instance_idx}.json") as f:
        key_to_logits = json.loads(f.read())
    for key, value in tid_dict.items(): 
        orig_img_path = f"structured_framework/visuals/flickr30k-clip-{instance_idx}-{0}-image.png"
        orig_text = open(f"structured_framework/visuals/flickr30k-clip-{instance_idx}-{0}-text.txt").read()
        doublegrad_img_path = f"structured_framework/visuals/flickr30k-clip-{key}-doublegrad.png"
        new_img_path = f"structured_framework/visuals/flickr30k-clip-{instance_idx}-0-saliency.png"
        # new_text = open(f"structured_framework/visuals/flickr30k-clip-{key}-new_text.txt").read()
        random_img_path =  f"structured_framework/visuals/flickr30k-clip-{key}-new_box_img_with_boxes.jpg"
        # gt_img_path = f"structured_framework/visuals/flickr30k-clip-{key}-gt_img.jpg"
        slide = prs.slides.add_slide(blank_slide_layout)


        # "original_logits": 10.043411,
        # "doublegrad_logits": 3.0287437,
        # "random_drop_logits": 8.995966,
        # "ground_truth_logits": 8.610129

        orig_lgts = key_to_logits[str(instance_idx)][key]["original_logits"]
        # dg_lgts = key_to_logits[str(instance_idx)][key]["doublegrad_box_logits"]
        # rd_lgts = key_to_logits[str(instance_idx)][key]["random_box_logits"]
        # gt_lgts = key_to_logits[str(instance_idx)][key]["ground_truth_logits"]


        # Original Image
        left = top = Inches(0.5)
        height = Inches(2.5)
        pic = slide.shapes.add_picture(orig_img_path, left, top, height=height)
        
        # Original Text
        left = Inches(1)
        top = Inches(3)
        height = Inches(2)
        width = Inches(3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        p = tf.add_paragraph()
        if len(orig_text.split()) > 15:
            third = len(orig_text.split()) // 3
            text = ' '.join(orig_text.split()[:third] + ["\n"]+ orig_text.split()[third:third*2] + ["\n"] + orig_text.split()[third*2:])
        else:
            mid = len(orig_text.split()) // 2
            text = ' '.join(orig_text.split()[:mid] + ["\n"] + orig_text.split()[mid:])
        p.text = text
        p.word_wrap = True
        p.font.size = Pt(10)

        left = Inches(1)
        top = Inches(3.5)
        height = Inches(0.5)
        width = Inches(3)
        txBox2 = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox2.text_frame

        double_grad_text = value["text"]

        if len(double_grad_text.split()) > 15:
            third = len(double_grad_text.split()) // 3
            text = ' '.join(double_grad_text.split()[:third] + ["\n"]+ double_grad_text.split()[third:third*2] + ["\n"] + double_grad_text.split()[third*2:])
        elif len(double_grad_text.split()) > 10:
            mid = len(orig_text.split()) // 2
            text = ' '.join(orig_text.split()[:mid] + ["\n"] + orig_text.split()[mid:])
        else:
            text = double_grad_text

        p = tf.add_paragraph()
        p.text = f"Double Grad Text: {text} \n Orig Logits: {orig_lgts:.3f}"
        p.word_wrap = True
        p.font.size = Pt(10)

        # DoubleGrad Image
        left = Inches(5)
        top = Inches(0.5)
        height = Inches(3)
        pic = slide.shapes.add_picture(doublegrad_img_path, left, top, height=height)

        # # New Text
        # left = Inches(5)
        # top = Inches(3)
        # height = Inches(2)
        # width = Inches(3)
        # txBox = slide.shapes.add_textbox(left, top, width, height)
        # tf = txBox.text_frame

        # p = tf.add_paragraph()
        # if len(new_text.split()) > 15:
        #     third = len(new_text.split()) // 3
        #     text = ' '.join(new_text.split()[:third] + ["\n"]+ new_text.split()[third:third*2] + ["\n"] + new_text.split()[third*2:])
        # else:
        #     mid = len(new_text.split()) // 2
        #     text = ' '.join(new_text.split()[:mid] + ["\n"] + new_text.split()[mid:])

        # p.text = text
        # p.word_wrap = True
        # p.font.size = Pt(10)

        # New Image
        left = Inches(0.5)
        top = Inches(4.5)
        height = Inches(3)
        pic = slide.shapes.add_picture(new_img_path, left, top, height=height)

        # New Image Logits
        left = Inches(0.5)
        top = Inches(4)
        height = Inches(1)
        width = Inches(3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
    
        p.text = f"Unimodal Gradient"
        p.word_wrap = True
        p.font.size = Pt(10)


        # # GT Image
        # left = Inches(3.75)
        # top = Inches(5)
        # height = Inches(2)
        # pic = slide.shapes.add_picture(gt_img_path, left, top, height=height)

        # # GT Image Logits
        # left = Inches(3.75)
        # top = Inches(4.5)
        # height = Inches(1)
        # width = Inches(3)
        # txBox = slide.shapes.add_textbox(left, top, width, height)
        # tf = txBox.text_frame
        # p = tf.add_paragraph()

        # p.text = f"Ground Truth Logits: {gt_lgts:.3f} (Delta: {gt_lgts - orig_lgts:.3f})"
        # p.word_wrap = True
        # p.font.size = Pt(10)

        # Random Image
        left = Inches(7)
        top = Inches(5)
        height = Inches(2)
        pic = slide.shapes.add_picture(random_img_path, left, top, height=height)

        # # Random Image Logits
        # left = Inches(7)
        # top = Inches(4.5)
        # height = Inches(1)
        # width = Inches(3)
        # txBox = slide.shapes.add_textbox(left, top, width, height)
        # tf = txBox.text_frame
        # p = tf.add_paragraph()

        # p.text = f"Random Image Logits: {rd_lgts:.3f} (Delta: {rd_lgts - orig_lgts:.3f})"
        # p.word_wrap = True
        # p.font.size = Pt(10)

prs.save('test.pptx')
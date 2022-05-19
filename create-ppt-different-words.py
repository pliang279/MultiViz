from pptx import Presentation, enum
from pptx.util import Inches, Pt

# CLIP Flickr30k Analysis - with Double Grad

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

target_ids_100 =  {0: '[CLS]', 1: 'a', 2: 'large', 3: 'bearded', 4: 'man', 5: 'flip', 6: '##s', 7: 'a', 8: 'cr', 9: '##ep', 10: '##e', 11: 'or', 12: 'om', 13: '##ele', 14: '##t', 15: 'in', 16: 'mid', 17: '##air', 18: 'with', 19: 'his', 20: 'fry', 21: '##ing', 22: 'pan', 23: '.', 24: '[SEP]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

instance_text_target_ids_100 = {
    "100_1": {"ids": [2], "text": "large"},
    "100_2": {"ids": [3], "text": "bearded"},
    "100_3": {"ids": [4], "text": "man"},
    "100_4": {"ids": [8,9,10], "text": "crepe"},
    "100_5": {"ids": [12,13,14], "text": "omelet"},
    "100_6": {"ids": [20, 21, 22], "text": "frying pan"},
    "100_7": {"ids": [2, 3, 4], "text": "large bearded man"},
    "100_8": {"ids": [2, 3, 4, 5, 6, 7, 8, 9, 10], "text": "large bearded man flips a crepe"},
    "100_9": {"ids": [2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12, 13, 14], "text": "large bearded man flips a crepe or omelet"},
    "100_10": {"ids": [12, 13, 14, 15, 16, 17], "text": "omelet in mid air"}
}

target_ids_150 = {0: '[CLS]', 1: 'a', 2: 'black', 3: 'dog', 4: 'with', 5: 'white', 6: 'facial', 7: 'and', 8: 'chest', 9: 'markings', 10: 'standing', 11: 'in', 12: 'chest', 13: 'high', 14: 'water', 15: '.', 16: '[SEP]', 17: '[PAD]', 18: '[PAD]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

instance_text_target_ids_150 = {
    "150_1": {"ids": [2], "text": "black"},
    "150_2": {"ids": [3], "text": "dog"},
    "150_3": {"ids": [5], "text": "white"},
    "150_4": {"ids": [6], "text": "facial"},
    "150_5": {"ids": [8], "text": "chest"},
    "150_6": {"ids": [9], "text": "markings"},
    "150_7": {"ids": [10], "text": "standing"},
    "150_8": {"ids": [11], "text": "in"},
    "150_9": {"ids": [12], "text": "chest"},
    "150_10": {"ids": [13], "text": "high"},
    "150_11": {"ids": [14], "text": "water"},
    "150_12": {"ids": [2, 3], "text": "black dog"},
    "150_13": {"ids":[5, 6], "text": "white facial"},
    "150_14": {"ids": [5, 6, 7, 8, 9], "text": "white facial and chest markings"},
    "150_15": {"ids": [12, 13, 14], "text": "chest high water"}
}

target_ids_200 = {0: '[CLS]', 1: 'a', 2: 'man', 3: 'is', 4: 'taking', 5: 'photographs', 6: 'of', 7: 'a', 8: 'large', 9: 'garden', 10: 'of', 11: 'white', 12: 'and', 13: 'orange', 14: 'tu', 15: '##lip', 16: '##s', 17: '.', 18: '[SEP]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

instance_text_target_ids_200 = {
    "200_1": {"ids": [2], "text": "man"},
    "200_2": {"ids": [5], "text": "photographs"},
    "200_3": {"ids": [9], "text": "garden"},
    "200_4": {"ids": [11], "text": "white"},
    "200_5": {"ids": [13], "text": "orange"},
    "200_6": {"ids": [14, 15, 16], "text": "tulips"},
    "200_7": {"ids": [1, 2, 3, 4, 5], "text": "a man is taking photographs"},
    "200_8": {"ids": [8, 9], "text": "large garden"},
    "200_9": {"ids": [1, 2, 3, 4, 5, 6, 7, 8, 9], "text": "a man is taking photographs of a large garden"},
    "200_10": {"ids": [8, 9, 10, 11, 12, 13, 14, 15, 16], "text": "a large garden of white and orange tulips"},
    "200_11": {"ids": [11, 12, 13, 14, 15, 16], "text": "white and orange tulips"},
    "200_12": {"ids": [1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12, 13, 14, 15, 16], "text": "a man is taking photographs of a large garden of white and orange tulips"}

}

target_ids_500 = {0: '[CLS]', 1: 'a', 2: 'little', 3: 'girl', 4: 'in', 5: 'front', 6: 'a', 7: 'pink', 8: 'food', 9: 'tray', 10: 'is', 11: 'getting', 12: 'her', 13: 'bike', 14: 'helmet', 15: 'on', 16: 'by', 17: 'a', 18: 'woman', 19: '.', 20: '[SEP]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

instance_text_target_ids_500 = {
    "500_1": {"ids": [2], "text": "little"},
    "500_2": {"ids": [3], "text": "girl"},
    "500_3": {"ids": [7], "text": "pink"},
    "500_4": {"ids": [8], "text": "food"},
    "500_5": {"ids": [9], "text": "tray"},
    "500_6": {"ids": [13], "text": "bike"},
    "500_7": {"ids": [14], "text": "helmet"},
    "500_8": {"ids": [18], "text": "woman"},
    "500_9": {"ids": [2, 3], "text": "little girl"},
    "500_10": {"ids": [8, 9], "text": "food tray"},
    "500_11": {"ids": [7, 8, 9], "text": "pink food tray"},
    "500_12": {"ids": [2, 3, 4, 5, 6, 7, 8, 9], "text": "little girl in front a pink food tray"},
    "500_13": {"ids": [13, 14], "text": "bike helmet"}
}

target_ids_50 = {0: '[CLS]', 1: 'three', 2: 'small', 3: 'dogs', 4: ',', 5: 'two', 6: 'white', 7: 'and', 8: 'one', 9: 'black', 10: 'and', 11: 'white', 12: ',', 13: 'on', 14: 'a', 15: 'sidewalk', 16: '.', 17: '[SEP]', 18: '[PAD]', 19: '[PAD]', 20: '[PAD]', 21: '[PAD]', 22: '[PAD]', 23: '[PAD]', 24: '[PAD]', 25: '[PAD]', 26: '[PAD]', 27: '[PAD]', 28: '[PAD]', 29: '[PAD]', 30: '[PAD]', 31: '[PAD]', 32: '[PAD]', 33: '[PAD]', 34: '[PAD]', 35: '[PAD]', 36: '[PAD]', 37: '[PAD]', 38: '[PAD]', 39: '[PAD]'}

instance_text_target_ids_50 = {
    "50_1": {"ids": [2], "text": "small"},
    "50_2": {"ids": [3], "text": "dogs"},
    "50_3": {"ids": [1], "text": "three"},
    "50_4": {"ids": [2, 3], "text": "small dogs"},
    "50_5": {"ids": [1, 2, 3], "text": "three small dogs"},
    "50_6": {"ids": [6], "text": "white"},
    "50_7": {"ids": [9], "text": "black"},
    "50_8": {"ids": [9, 10, 11], "text": "black and white"},
    "50_9": {"ids": [15], "text": "sidewalk"},
    "50_10": {"ids": [5, 6, 7, 8, 9, 10, 11], "text": "two white and one black and white"},
    "50_11": {"ids": [1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11], "text": "three small dogs, two white and one black and white"},
}


id_to_tids = {
    50: instance_text_target_ids_50,
    100: instance_text_target_ids_100,
    150: instance_text_target_ids_150,
    200: instance_text_target_ids_200,
    500: instance_text_target_ids_500
}

for instance_idx, tid_dict in id_to_tids.items():
    for key, value in tid_dict.items(): 
        orig_img_path = f"structured-framework/visuals/flickr30k-vilt-{instance_idx}-{0}-image.png"
        orig_text = open(f"structured-framework/visuals/flickr30k-clip-{instance_idx}-{0}-text.txt").read()
        doublegrad_img_path = f"structured-framework/visuals/flickr30k-vilt-{key}-doublegrad.png"
        slide = prs.slides.add_slide(blank_slide_layout)

        # Original Image
        left = top = Inches(0.5)
        height = Inches(2.5)
        pic = slide.shapes.add_picture(orig_img_path, left, top, height=height)
        
        # Original Text
        left = Inches(1)
        top = Inches(3)
        height = Inches(2)
        width = Inches(3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        p = tf.add_paragraph()
        if len(orig_text.split()) > 15:
            third = len(orig_text.split()) // 3
            text = ' '.join(orig_text.split()[:third] + ["\n"]+ orig_text.split()[third:third*2] + ["\n"] + orig_text.split()[third*2:])
        else:
            mid = len(orig_text.split()) // 2
            text = ' '.join(orig_text.split()[:mid] + ["\n"] + orig_text.split()[mid:])
        p.text = text
        p.word_wrap = True
        p.font.size = Pt(10)

        left = Inches(1)
        top = Inches(3.5)
        height = Inches(0.5)
        width = Inches(3)
        txBox2 = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox2.text_frame

        double_grad_text = value["text"]

        if len(double_grad_text.split()) > 15:
            third = len(double_grad_text.split()) // 3
            text = ' '.join(double_grad_text.split()[:third] + ["\n"]+ double_grad_text.split()[third:third*2] + ["\n"] + double_grad_text.split()[third*2:])
        elif len(double_grad_text.split()) > 10:
            mid = len(orig_text.split()) // 2
            text = ' '.join(orig_text.split()[:mid] + ["\n"] + orig_text.split()[mid:])
        else:
            text = double_grad_text

        p = tf.add_paragraph()
        p.text = "Double Grad Text: " + text
        p.word_wrap = True
        p.font.size = Pt(10)


        # DoubleGrad Image
        left = Inches(5)
        top = Inches(0)
        height = Inches(3.5)
        pic = slide.shapes.add_picture(doublegrad_img_path, left, top, height=height)

prs.save('test.pptx')